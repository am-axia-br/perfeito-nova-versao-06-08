# ===================================================================
# app/rag_manager.py (OTIMIZADO COM EMBEDDINGS JURÍDICOS EM PORTUGUÊS)
#
# Mudanças:
# - Implementação de embeddings especializados para textos jurídicos
# - Configurado para usar juris-bert-base-portuguese como principal
# - Fallback para Google Embeddings se necessário
# - Mantida toda a funcionalidade existente
# ===================================================================

import os
import io
import chromadb
import fitz
import pandas as pd
import csv
import json
import datetime

from odf.opendocument import load
from odf.text import P, H
from bs4 import BeautifulSoup
from pptx import Presentation
from docx import Document

from langchain.text_splitter import RecursiveCharacterTextSplitter
from langchain_google_genai import GoogleGenerativeAIEmbeddings

# Novos imports para embeddings especializados
from sentence_transformers import SentenceTransformer
import torch

# --- CONFIGURAÇÕES E CONSTANTES ---
CHROMA_HOST = os.getenv("CHROMA_HOST", "localhost")
CHROMA_PORT = int(os.getenv("CHROMA_PORT", 8000))
COLLECTION_NAME = "documentos_juridicos"

# --- SINGLETONS ---
_chroma_client = None
_embedding_function = None

# --- ADAPTADOR DE EMBEDDING ATUALIZADO ---
class EmbeddingFunctionAdapter:
    def __init__(self, model_name="neuralmind/juris-bert-base-portuguese"):
        """
        Adaptador para modelos de embedding do HuggingFace
        
        Args:
            model_name: Nome do modelo no HuggingFace
                - "BAAI/bge-base-en-v1.5" (multilíngue geral)
                - "neuralmind/juris-bert-base-portuguese" (português jurídico)
        """
        print(f"🔄 Carregando modelo de embedding: {model_name}")
        try:
            # Carrega o modelo localmente usando sentence-transformers
            self.model = SentenceTransformer(model_name)
            self.device = "cuda" if torch.cuda.is_available() else "cpu"
            self.model.to(self.device)
            print(f"✅ Modelo de embedding carregado com sucesso (dispositivo: {self.device})!")
        except Exception as e:
            print(f"❌ Erro ao carregar modelo de embedding: {e}")
            import traceback
            traceback.print_exc()
            raise
    
    def __call__(self, input):
        """Interface compatível com ChromaDB"""
        if not input or len(input) == 0:
            print("⚠️ Lista de entrada vazia para embeddings")
            return []
        
        try:
            print(f"🔢 Gerando embeddings para {len(input)} texto(s)")
            
            # Transforma os textos em embeddings
            embeddings = self.model.encode(input, convert_to_numpy=True).tolist()
            
            if len(embeddings) != len(input):
                print(f"⚠️ Número de embeddings ({len(embeddings)}) não corresponde à entrada ({len(input)})")
            
            return embeddings
            
        except Exception as e:
            print(f"❌ Erro ao gerar embeddings: {e}")
            import traceback
            traceback.print_exc()
            return [[] for _ in input]

# --- ADAPTADOR PARA COMPATIBILIDADE COM LANGCHAIN (usado no fallback) ---
class LangchainEmbeddingAdapter:
    def __init__(self, langchain_embeddings):
        self.langchain_embeddings = langchain_embeddings
        print("✅ LangchainEmbeddingAdapter criado com sucesso")
    
    def __call__(self, input):
        """Interface compatível com ChromaDB 0.4.16+"""
        # Tratamento de entrada vazia
        if not input or len(input) == 0:
            print("⚠️ Aviso: Lista de entrada vazia para embeddings")
            # Retorna lista vazia em vez de falhar
            return []
        
        try:
            # Para debugging - mostra quantidade de textos
            print(f"🔢 Gerando embeddings para {len(input)} texto(s)")
            
            # Usa o método embed_documents do LangChain
            embeddings = self.langchain_embeddings.embed_documents(input)
            
            # Verificação dos resultados
            if len(embeddings) != len(input):
                print(f"⚠️ Aviso: Número de embeddings ({len(embeddings)}) não corresponde à entrada ({len(input)})")
            
            return embeddings
            
        except Exception as e:
            print(f"❌ Erro ao gerar embeddings: {e}")
            import traceback
            traceback.print_exc()
            # Retorna embeddings vazios para evitar quebra total
            return [[] for _ in input]

def get_embedding_function():
    global _embedding_function
    if _embedding_function is None:
        try:
            # Opção 1: Modelo multilíngue geral (bom para vários idiomas)
            # _embedding_function = EmbeddingFunctionAdapter("BAAI/bge-base-en-v1.5")
            
            # Opção 2: Modelo especializado em português jurídico (recomendado)
            _embedding_function = EmbeddingFunctionAdapter("neuralmind/juris-bert-base-portuguese")
            
            print("✅ Função de embedding jurídico em português criada com sucesso!")
        except Exception as e:
            print(f"❌ Erro ao criar embedding function: {e}")
            # Fallback para o Google Embeddings caso falhe
            print("⚠️ Usando Google Embeddings como fallback")
            from auth_manager import configure_gemini_auth
            import os
            
            AUTH_METHOD = configure_gemini_auth()
            from langchain_google_genai import GoogleGenerativeAIEmbeddings
            
            if AUTH_METHOD == "service_account":
                langchain_embeddings = GoogleGenerativeAIEmbeddings(
                    model="models/embedding-001"
                )
            elif AUTH_METHOD == "api_key":
                langchain_embeddings = GoogleGenerativeAIEmbeddings(
                    model="models/embedding-001",
                    google_api_key=os.getenv("GEMINI_API_KEY")
                )
            else:
                raise RuntimeError("Autenticação Gemini não configurada corretamente.")
            
            # Usamos o adaptador definido acima em vez de importar
            _embedding_function = LangchainEmbeddingAdapter(langchain_embeddings)
            
            import traceback
            traceback.print_exc()
    
    return _embedding_function

# --- CLIENTE CHROMADB (SEM TENTATIVAS) ---
def get_chroma_client():
    global _chroma_client
    if _chroma_client is not None:
        return _chroma_client

    try:
        client = chromadb.HttpClient(host=CHROMA_HOST, port=CHROMA_PORT)
        client.heartbeat()
        print(f"✅ Conectado ao ChromaDB em http://{CHROMA_HOST}:{CHROMA_PORT}")
        _chroma_client = client
        return _chroma_client
    except Exception as e:
        print(f"❌ Erro ao conectar ao ChromaDB: {e}")
        return None

# --- COLEÇÃO ---

def get_or_create_collection():
    try:
        client = get_chroma_client()
        if not client:
            print("❌ ChromaDB indisponível - cliente é None")
            return None
        
        embedding_func = get_embedding_function()
        
        # Primeiro verificamos se a coleção já existe
        collections = client.list_collections()
        exists = any(col.name == COLLECTION_NAME for col in collections)
        
        if exists:
            print(f"✅ Coleção {COLLECTION_NAME} encontrada, recuperando...")
            return client.get_collection(name=COLLECTION_NAME, embedding_function=embedding_func)
        else:
            print(f"🆕 Criando nova coleção {COLLECTION_NAME}...")
            return client.create_collection(
                name=COLLECTION_NAME,
                embedding_function=embedding_func,
                metadata={"hnsw:space": "cosine"}
            )
    except Exception as e:
        print(f"❌ Erro ao criar/obter coleção: {e}")
        import traceback
        traceback.print_exc()
        return None

# --- STATUS ---
def get_rag_status():
    client = get_chroma_client()
    if not client:
        return {"connected": False, "doc_count": 0, "message": "ChromaDB indisponível."}
    try:
        collection = client.get_or_create_collection(name=COLLECTION_NAME)
        count = collection.count()
        msg = "Base de conhecimento ativa." if count > 0 else "Base conectada, mas vazia."
        return {"connected": True, "doc_count": count, "message": msg}
    except Exception as e:
        return {"connected": False, "doc_count": 0, "message": str(e)}

# --- EXTRAÇÃO DE TEXTO ---
def _extrair_texto_de_arquivo(file_bytes, file_name):
    """
    Extrai texto de diferentes formatos de arquivo com tratamento de erros aprimorado.
    """
    try:
        if file_name.endswith(".pdf"):
            with fitz.open(stream=file_bytes, filetype="pdf") as doc:
                return "".join(page.get_text() for page in doc)

        elif file_name.endswith(".docx"):
            doc = Document(io.BytesIO(file_bytes))
            return "\n".join(p.text for p in doc.paragraphs)

        elif file_name.endswith(".txt"):
            return file_bytes.decode("utf-8")

        elif file_name.endswith(".xlsx"):
            try:
                # Verificar se openpyxl está disponível
                import openpyxl
                
                xls = pd.ExcelFile(io.BytesIO(file_bytes))
                partes = []
                for aba in xls.sheet_names:
                    df = pd.read_excel(xls, sheet_name=aba)
                    partes.append(f"📄 Aba: {aba}\n{df.to_string(index=False)}")
                return "\n\n".join(partes)
            except ImportError:
                mensagem = "❌ Biblioteca 'openpyxl' não instalada. Execute 'pip install openpyxl' e tente novamente."
                print(mensagem)
                # Retornar uma mensagem em vez de string vazia para feedback ao usuário
                return mensagem
            except Exception as excel_error:
                print(f"❌ Erro ao processar arquivo Excel '{file_name}': {excel_error}")
                return f"❌ Erro ao processar Excel: {excel_error}"

        elif file_name.endswith(".pptx"):
            prs = Presentation(io.BytesIO(file_bytes))
            slides = []
            for i, slide in enumerate(prs.slides):
                textos = [shape.text for shape in slide.shapes if hasattr(shape, "text")]
                slides.append(f"📽️ Slide {i+1}:\n" + "\n".join(textos))
            return "\n\n".join(slides)

        elif file_name.endswith(".html"):
            soup = BeautifulSoup(file_bytes, "html.parser")
            return soup.get_text(separator="\n")

        elif file_name.endswith(".csv"):
            decoded = file_bytes.decode('utf-8').splitlines()
            reader = csv.reader(decoded)
            return "\n".join([", ".join(row) for row in reader])

        elif file_name.endswith(".json"):
            decoded = file_bytes.decode('utf-8')
            json_data = json.loads(decoded)
            return json.dumps(json_data, indent=2, ensure_ascii=False)

        elif file_name.endswith(".odt"):
            odt = load(io.BytesIO(file_bytes))
            paragraphs = odt.getElementsByType(P) + odt.getElementsByType(H)
            return "\n".join([str(p.firstChild.data) for p in paragraphs if p.firstChild])

        else:
            raise ValueError(f"❌ Formato não suportado: {file_name}")

    except Exception as e:
        print(f"❌ Erro ao extrair texto de {file_name}: {e}")
        import traceback
        traceback.print_exc()  # Imprime o stack trace completo para debugging
        return f"❌ Erro ao processar o arquivo: {str(e)}"
# --- ADICIONAR DOCUMENTO ---

def adicionar_documento_ao_rag(file_bytes, file_name):
    try:
        collection = get_or_create_collection()
        if not collection:
            return f"❌ ChromaDB indisponível."

        # Verificar se o arquivo já foi processado recentemente (nos últimos 5 minutos)
        log_path = "logs/treinamentos.jsonl"
        if os.path.exists(log_path):
            try:
                with open(log_path, encoding="utf-8") as f:
                    linhas = f.readlines()
                    for linha in reversed(linhas[-10:]):  # Verifica as últimas 10 entradas
                        registro = json.loads(linha)
                        if registro["arquivo"] == file_name:
                            # Calcula quanto tempo passou desde o último processamento
                            ultima_data = datetime.datetime.fromisoformat(registro["data"])
                            tempo_passado = datetime.datetime.now() - ultima_data
                            # Se foi processado nos últimos 5 minutos, ignora
                            if tempo_passado.total_seconds() < 300:  # 5 minutos
                                print(f"ℹ️ Ignorando reprocessamento de '{file_name}' (processado há {tempo_passado.total_seconds():.0f} segundos)")
                                return f"ℹ️ Arquivo '{file_name}' já foi processado recentemente."
            except Exception as e:
                print(f"⚠️ Erro ao verificar log de processamento: {e}")

        # Limpar documentos existentes com o mesmo nome para evitar duplicação
        try:
            existentes = collection.get(where={"source": file_name})
            if existentes and existentes['ids']:
                print(f"🧹 Removendo {len(existentes['ids'])} chunks existentes do arquivo {file_name}")
                collection.delete(ids=existentes['ids'])
        except Exception as e:
            print(f"⚠️ Não foi possível limpar documentos existentes: {e}")

        texto = _extrair_texto_de_arquivo(file_bytes, file_name)
        if not texto.strip():
            return f"⚠️ Arquivo '{file_name}' está vazio ou ilegível."

        # Verificar se o texto tem conteúdo significativo
        if texto.startswith("❌ Erro") or texto.startswith("❌ Biblioteca"):
            print(f"⚠️ Erro ao extrair texto de '{file_name}': {texto}")
            return texto  # Retorna a mensagem de erro que veio da extração

        splitter = RecursiveCharacterTextSplitter(chunk_size=1500, chunk_overlap=200)
        chunks = splitter.split_text(texto)
        if not chunks:
            return f"⚠️ Nenhum trecho extraído de '{file_name}'."

        # Registrar que está iniciando o processamento
        print(f"📝 Processando {len(chunks)} chunks para '{file_name}'")
        
        # Gerar IDs únicos com timestamp para evitar conflitos
        timestamp = int(datetime.datetime.now().timestamp())
        
        for i in range(0, len(chunks), 100):
            fatia = chunks[i:i + 100]
            ids = [f"{file_name}-{timestamp}-{i + j}" for j in range(len(fatia))]
            metadatas = [{"source": file_name, "timestamp": timestamp} for _ in fatia]
            collection.add(documents=fatia, metadatas=metadatas, ids=ids)

        registro = {
            "arquivo": file_name,
            "data": datetime.datetime.now().isoformat(),
            "quantidade_chunks": len(chunks),
            "timestamp": timestamp
        }
        log_path = "logs/treinamentos.jsonl"
        os.makedirs(os.path.dirname(log_path), exist_ok=True)
        with open(log_path, "a", encoding="utf-8") as f:
            f.write(json.dumps(registro, ensure_ascii=False) + "\n")

        return f"✅ Treinamento concluído com {file_name}. Adicionados {len(chunks)} chunks."

    except Exception as e:
        print(f"❌ Erro ao adicionar '{file_name}' ao RAG: {e}")
        import traceback
        traceback.print_exc()
        return f"❌ Erro ao treinar com {file_name}: {str(e)}"

# --- CONSULTAR ---
def consultar_rag(query_text, n_results=5):
    try:
        collection = get_or_create_collection()
        if not collection:
            print("⚠️ ChromaDB indisponível.")
            return ""

        print(f"🔍 Consultando ChromaDB com: '{query_text[:50]}...'")
        results = collection.query(query_texts=[query_text], n_results=n_results)
        
        # Logging detalhado
        documents = results.get("documents", [[]])
        num_docs = len(documents[0]) if documents and documents[0] else 0
        print(f"📊 Consulta RAG: recuperados {num_docs} documentos")
        
        documentos = results.get("documents", [[]])[0]
        if not documentos:
            print("⚠️ Nenhum resultado encontrado.")
            return ""

        contexto = "\n\n---\n\n".join(documentos)
        print("✅ Contexto recuperado com sucesso.")
        return contexto

    except Exception as e:
        print(f"❌ Erro ao consultar o RAG: {e}")
        import traceback
        traceback.print_exc()
        return ""

# --- ZERAR BASE DE CONHECIMENTO ---
def zerar_base_rag():
    """
    Limpa completamente a base de conhecimento do RAG.
    Útil para fins de teste e reinicialização.
    """
    try:
        client = get_chroma_client()
        if not client:
            return {"success": False, "message": "❌ ChromaDB indisponível."}
        
        # Verifica se a coleção existe
        collections = client.list_collections()
        exists = any(col.name == COLLECTION_NAME for col in collections)
        
        if exists:
            # Deleta a coleção existente
            client.delete_collection(name=COLLECTION_NAME)
            print(f"🧹 Coleção {COLLECTION_NAME} removida com sucesso")
        
        # Recria uma coleção vazia
        embedding_func = get_embedding_function()
        client.create_collection(
            name=COLLECTION_NAME,
            embedding_function=embedding_func,
            metadata={"hnsw:space": "cosine"}
        )
        print(f"🆕 Coleção {COLLECTION_NAME} recriada vazia")
        
        # Limpa o log de treinamentos
        log_path = "logs/treinamentos.jsonl"
        if os.path.exists(log_path):
            os.remove(log_path)
            os.makedirs(os.path.dirname(log_path), exist_ok=True)
            with open(log_path, "w", encoding="utf-8") as f:
                pass  # Cria um arquivo vazio
            print("📝 Log de treinamentos limpo")
        
        # Limpa o singleton do cliente para forçar reconexão
        global _chroma_client
        _chroma_client = None
        
        return {"success": True, "message": "✅ Base de conhecimento zerada com sucesso!"}
    
    except Exception as e:
        print(f"❌ Erro ao zerar a base RAG: {e}")
        import traceback
        traceback.print_exc()
        return {"success": False, "message": f"❌ Erro ao zerar base: {e}"}
    
def get_chroma_content(filter_source=None, limit=1000):
    """
    Recupera o conteúdo do ChromaDB para visualização em tabela.
    
    Args:
        filter_source: Filtrar por nome do arquivo fonte (opcional)
        limit: Limite máximo de registros
    
    Returns:
        Dicionário com os dados e estatísticas
    """
    try:
        collection = get_or_create_collection()
        if not collection:
            return {"success": False, "message": "ChromaDB indisponível"}
        
        # Aplicar filtro se especificado
        where_filter = {"source": filter_source} if filter_source else None
        
        # Recuperar documentos
        result = collection.get(
            where=where_filter,
            limit=limit
        )
        
        # Contar documentos por fonte
        sources = {}
        for meta in result.get("metadatas", []):
            if meta and "source" in meta:
                source = meta["source"]
                if source not in sources:
                    sources[source] = 0
                sources[source] += 1
        
        return {
            "success": True,
            "data": result,
            "count": len(result.get("ids", [])),
            "sources": sources
        }
    except Exception as e:
        print(f"❌ Erro ao recuperar dados do ChromaDB: {e}")
        import traceback
        traceback.print_exc()
        return {"success": False, "message": str(e)}